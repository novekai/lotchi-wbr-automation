#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
wbr_write_text.py — pose les textes rédigés par Claude dans les cellules du deck WBR.

Séparation des rôles : le script place le texte, Claude fournit le jugement (rédaction).
Aucune donnée n'est calculée ici ; on écrit seulement ce que le JSON contient.

Cible : le PPTX déjà produit par generate_wbr.py. On écrit dans :
  - km_ap_table_s5 / s7 / s8 / s9 / s11  -> Key Metric (col 0) + Action Plan (col 1).
    Villes internationales : la colonne 2 de la s11 s'intitule « Reviews » et arrive PRÉ-REMPLIE
    (notes Fever/Google) -> on n'écrit la colonne 2 que sous un en-tête « Action Plan ».
  - actionplan_table (« To do for next week »)  -> une ligne par action (Lotchi / Fever /
    Partner facultatif / Deadline).
  - toute table « *_comments_table » (slides campagne) -> colonne « Comments and recommendations »,
    appariée par le nom de campagne de la 1re colonne, avec préfixe de type facultatif
    « TYPE:nom » (TYPE = ce qui suit « Campaign review: » dans le titre de la slide) pour
    distinguer un même nom présent sur plusieurs slides.

On ne touche JAMAIS à checklist_w1_table (Plan d'action W-1, laissé vide), ni aux graphiques,
ni aux tables de chiffres (bigpicture_*, *_kpi_summary_table, sold-out).

INTÉGRITÉ DU CONTENEUR : avant d'écrire, on supprime les slides « fantômes » (parts reliées mais
hors sommaire, laissées par generate_wbr.py) ; après sauvegarde, on vérifie qu'aucun nom de part
n'est dupliqué. Sans ça, PowerPoint refuse d'ouvrir le fichier.

Usage :
  python wbr/wbr_write_text.py --deck WBR_Bayonne_W29.pptx --texts textes.json [--out WBR_final.pptx]

Schéma de textes.json :
{
  "key_metrics": {
    "5":  {"key_metric": "...", "action_plan": "..."},
    "7":  {"key_metric": "...", "action_plan": "N/A"},
    "8":  {...}, "9": {...}, "11": {...}
  },
  "todo": [
    {"action": "Baisser le spend", "lotchi": true, "fever": false, "partner": false, "deadline": ""},
    {"action": "Transmettre les recos Google", "lotchi": false, "fever": true, "deadline": "W30"}
  ],
  "campaign_comments": {
    "AC - Influencers": "Commentaire pour un nom unique dans le deck.",
    "REACH:AC - Nocturnes": "Commentaire spécifique à la slide REACH.",
    "CONVERSION:AC - Nocturnes": "Commentaire spécifique à la slide CONVERSION."
  }
}
Toutes les clés sont optionnelles : on ne remplit que ce qui est fourni. Pour campaign_comments,
utiliser le nom de campagne EXACT (1re colonne de la table) comme clé ; si un même nom apparaît
sur plusieurs slides, préfixer par le type lu dans le titre de la slide (« TYPE:nom »).

Dépendances : python-pptx.
"""
import argparse, json, sys, copy, zipfile
from collections import Counter
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

KM_SLIDES = ("5", "7", "8", "9", "11")
MARK = "X"  # marqueur dans les colonnes Lotchi / Fever du To do


# ── écriture d'une cellule en préservant les retours à la ligne ───────────────
def set_cell(cell, text, size=9, bold=False, align=None, color=None):
    text = "" if text is None else str(text)
    lines = text.split("\n")
    tf = cell.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        p.font.size = Pt(size)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    cell.margin_left = Pt(3)
    cell.margin_right = Pt(3)
    cell.margin_top = Pt(1)
    cell.margin_bottom = Pt(1)


def find_tables(prs):
    """Itère (slide_index, shape, table) sur toutes les tables du deck."""
    for si, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.has_table:
                yield si, sh, sh.table


def header(table):
    return [table.cell(0, c).text.strip() for c in range(len(table.columns))]


# ── intégrité du conteneur PPTX ───────────────────────────────────────────────
def drop_orphan_slides(prs):
    """Supprime les parts de slides orphelines : reliées à la présentation mais absentes du
    sommaire (sldIdLst). generate_wbr.py (reorder_campaign) retire les slides-modèles de campagne
    du sommaire sans supprimer leurs parts -> au ré-enregistrement, python-pptx réutilise ces noms
    de parts et produit des doublons (slide14.xml deux fois) -> PPTX que PowerPoint refuse d'ouvrir.
    On coupe la relation : la part devient inatteignable et n'est plus sérialisée."""
    live = {s.part for s in prs.slides}
    dropped = 0
    for rId, rel in list(prs.part.rels.items()):
        if rel.reltype == RT.SLIDE and rel.target_part not in live:
            prs.part.drop_rel(rId)
            dropped += 1
    return dropped


def assert_integrity(path):
    """Garde-fou final : aucun nom de part dupliqué dans le conteneur ZIP. Lève une erreur plutôt
    que de livrer un fichier corrompu silencieusement."""
    names = zipfile.ZipFile(path).namelist()
    dups = sorted(n for n, c in Counter(names).items() if c > 1)
    if dups:
        raise SystemExit(f"ECHEC integrite : parts dupliquees dans {path} -> {dups}")


# ── Key Metric / Action Plan ──────────────────────────────────────────────────
def fill_key_metrics(prs, km):
    done = []
    for _, sh, t in find_tables(prs):
        name = sh.name
        for s in KM_SLIDES:
            if name == f"km_ap_table_s{s}" and s in km:
                entry = km[s] or {}
                set_cell(t.cell(1, 0), entry.get("key_metric"), size=9)
                # Villes internationales : la col. 2 de la s11 s'intitule « Reviews » et arrive
                # PRÉ-REMPLIE (notes Fever/Google) — n'écrire que sous un en-tête « Action Plan ».
                if t.cell(0, 1).text.strip().lower() == "action plan":
                    set_cell(t.cell(1, 1), entry.get("action_plan"), size=9)
                done.append(s)
    return done


# ── To do for next week (actionplan_table) ────────────────────────────────────
def _add_row(table):
    """Ajoute une ligne en clonant la dernière <a:tr> puis en vidant ses cellules."""
    tbl = table._tbl
    new_tr = copy.deepcopy(tbl.tr_lst[-1])
    tbl.append(new_tr)
    r = len(table.rows) - 1
    for c in range(len(table.columns)):
        set_cell(table.cell(r, c), "")
    return r


def fill_todo(prs, todo):
    if not todo:
        return 0
    for _, sh, t in find_tables(prs):
        if sh.name != "actionplan_table":
            continue
        # colonnes : Actions | To be done by Lotchi | To be done by Fever | To be done by Partner | Deadline
        # Les positions varient selon le template : Lotchi/Fever/Partner sont repérées par
        # l'en-tête, Deadline reste la dernière colonne.
        heads = [c.text.strip().lower() for c in t.rows[0].cells]

        def col_idx(word, default):
            return next((i for i, h in enumerate(heads) if word in h), default)

        c_lotchi, c_fever = col_idx("lotchi", 1), col_idx("fever", 2)
        c_partner = col_idx("partner", None)
        deadline_col = len(t.columns) - 1
        for i, item in enumerate(todo):
            r = 1 if i == 0 else _add_row(t)
            set_cell(t.cell(r, 0), item.get("action"), size=9)
            set_cell(t.cell(r, c_lotchi), MARK if item.get("lotchi") else "", size=9, align=PP_ALIGN.CENTER)
            set_cell(t.cell(r, c_fever), MARK if item.get("fever") else "", size=9, align=PP_ALIGN.CENTER)
            if c_partner is not None:
                set_cell(t.cell(r, c_partner), MARK if item.get("partner") else "", size=9, align=PP_ALIGN.CENTER)
            set_cell(t.cell(r, deadline_col), item.get("deadline") or "", size=9, align=PP_ALIGN.CENTER)
        return len(todo)
    return 0


# ── Commentaires de campagne (appariés par nom de campagne) ───────────────────
def slide_campaign_type(slide):
    """Type de campagne d'une slide « Campaign review: X. », lu dans son titre (ex. REACH,
    CONVERSION, ENGAGEMENT ou tout futur objectif). '' si aucun titre reconnaissable — le deck
    fait autorité, aucune liste de types n'est figée ici."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt.lower().startswith("campaign review"):
                return txt.split(":", 1)[-1].strip().rstrip(".").upper()
    return ""


def fill_campaign_comments(prs, comments):
    if not comments:
        return 0
    n = 0
    for sl in prs.slides:
        ctype = None  # titre de slide lu au premier besoin seulement
        for sh in sl.shapes:
            if not sh.has_table:
                continue
            t = sh.table
            h = header(t)
            if not h or h[-1].lower() != "comments and recommendations":
                continue
            camp = t.cell(1, 0).text.strip()
            if not camp:
                continue
            if ctype is None:
                ctype = slide_campaign_type(sl)
            # 1) clé « TYPE:nom » (prioritaire : distingue un même nom de campagne présent
            #    sur plusieurs slides, ex. en REACH et en CONVERSION).
            match = comments.get(f"{ctype}:{camp}") if ctype else None
            # 2) correspondance EXACTE par nom seul (cas nominal, noms uniques).
            if match is None:
                match = comments.get(camp)
            # 3) sinon inclusion, mais seulement si NON AMBIGUË (un unique candidat).
            if match is None:
                cands = [v for k, v in comments.items() if camp in k or k in camp]
                if len(cands) == 1:
                    match = cands[0]
            if match is not None:
                set_cell(t.cell(1, len(t.columns) - 1), match, size=8)
                n += 1
    return n


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--deck", required=True, help="PPTX produit par generate_wbr.py")
    ap.add_argument("--texts", required=True, help="JSON des textes rédigés")
    ap.add_argument("--out", default=None, help="PPTX de sortie (défaut : écrase --deck)")
    a = ap.parse_args()

    with open(a.texts, encoding="utf-8") as f:
        texts = json.load(f)
    prs = Presentation(a.deck)

    orphans = drop_orphan_slides(prs)   # nettoie les slides fantômes AVANT toute écriture

    km = fill_key_metrics(prs, texts.get("key_metrics", {}))
    todo = fill_todo(prs, texts.get("todo", []))
    comm = fill_campaign_comments(prs, texts.get("campaign_comments", {}))

    out = a.out or a.deck
    prs.save(out)
    assert_integrity(out)               # refuse de livrer un PPTX au conteneur invalide
    print(f"OK -> {out} | key_metrics={sorted(km)} todo_rows={todo} "
          f"campaign_comments={comm} orphelines_supprimees={orphans}")


if __name__ == "__main__":
    main()
