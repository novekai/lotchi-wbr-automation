"""
Utilitaires pour éditer le template `wbr_template.pptx`.

Fonctions principales :
- find_by_name(slide, name)              : récupère un shape par son nom
- find_by_name_at_pos(slide, name, ...)  : désambiguïse si plusieurs shapes partagent un nom
- replace_picture(slide, shape, path)    : remplace une image en préservant position et taille
- set_cell_text(cell, text, ...)         : remplit une cellule de tableau avec une liste de lignes
- clear_table_keep_headers(table, n=1)   : vide toutes les lignes sous l'en-tête
- substitute_paragraphs(shape, paras)    : remplace les paragraphes d'une zone de texte ligne par ligne
- fit_all_tables(prs)                    : garde-fou anti-débordement des tableaux (à appeler avant save)
- dedupe_zip(pptx_path)                  : supprime les doublons XML du PPTX (essentiel après édition)
"""
import math  # wrap estimation
import shutil
import zipfile
from pathlib import Path

from pptx.util import Pt
from pptx.dml.color import RGBColor


# ---------- Recherche de shapes ----------

def find_by_name(slide, name):
    """Retourne le premier shape de la slide avec ce nom, ou None."""
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def find_all_by_name(slide, name):
    """Tous les shapes avec ce nom."""
    return [sh for sh in slide.shapes if sh.name == name]


def find_by_name_at_pos(slide, name, *, near_left=None, near_top=None, tol=0.5):
    """
    Récupère un shape par nom, filtré par position approximative (en inches).
    Utile quand plusieurs shapes partagent le même nom (cas Image 2 sur la slide 11
    avant renommage). Avec un template propre, ce cas ne devrait plus arriver,
    mais l'utilitaire est conservé pour robustesse.
    """
    EMU = 914400  # 1 inch
    candidates = find_all_by_name(slide, name)
    if not candidates:
        return None
    if near_left is None and near_top is None:
        return candidates[0]
    for sh in candidates:
        L = (sh.left or 0) / EMU
        T = (sh.top or 0) / EMU
        if near_left is not None and abs(L - near_left) > tol:
            continue
        if near_top is not None and abs(T - near_top) > tol:
            continue
        return sh
    return candidates[0]


# ---------- Édition d'images ----------

def replace_picture(slide, shape, new_path):
    """
    Remplace une image en préservant position, taille et z-order approximatif.
    Si `shape` est None, ne fait rien et retourne False.
    """
    if shape is None:
        return False
    L, T, W, H = shape.left, shape.top, shape.width, shape.height
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(new_path), L, T, width=W, height=H)
    return True


# ---------- Édition de tableaux ----------

def set_cell_text(cell, text, *, size=9, bold=False):
    """
    Remplit une cellule de tableau. `text` peut être :
    - une string : un seul paragraphe
    - une liste de strings : un paragraphe par élément (= une ligne par item)
    """
    # vider la cellule sans toucher au style de la cellule
    for p in list(cell.text_frame.paragraphs):
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
    if isinstance(text, str):
        text = [text]
    tf = cell.text_frame
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(0xE6, 0xE6, 0xE6)  # texte clair (deck dark)


def clear_table_keep_headers(table, n_keep_header_rows=1):
    """Vide toutes les cellules sous les `n_keep_header_rows` premières lignes."""
    for i in range(n_keep_header_rows, len(table.rows)):
        for j in range(len(table.columns)):
            cell = table.cell(i, j)
            for p in list(cell.text_frame.paragraphs):
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)


# ---------- Édition de paragraphes (zone de texte multi-lignes) ----------

def substitute_paragraphs(shape, paragraphs):
    """
    Remplace ligne par ligne les paragraphes d'une zone de texte, en préservant
    la mise en forme du premier run de chaque paragraphe.

    Utilisation typique : la cover_textbox qui contient 4 paragraphes statiques
    (« Business Review » / VILLE / DATE / « Version Claude »).
    """
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    existing = list(tf.paragraphs)
    for i, line in enumerate(paragraphs):
        if i >= len(existing):
            break
        runs = list(existing[i].runs)
        if not runs:
            continue
        runs[0].text = line
        for r in runs[1:]:
            r._r.getparent().remove(r._r)


# ---------- Garde-fou anti-débordement des tableaux ----------
#
# PowerPoint agrandit automatiquement les lignes d'un tableau quand le texte
# wrappe : la hauteur déclarée dans le XML n'est qu'un minimum. Un texte long
# injecté dans une table placée en bas de slide (ex. km_ap_table_s13) fait donc
# déborder la table hors de la page. Ces fonctions estiment la hauteur RENDUE
# de chaque table et, en cas de dépassement : réduisent la police des lignes de
# données (10 → 9 → 8 pt), puis remontent la table si nécessaire.

EMU_PER_IN = 914400
CHAR_W_RATIO = 0.5    # largeur moyenne d'un caractère ≈ 0.5 × taille de police
LINE_SPACING = 1.2    # interligne PowerPoint par défaut


def _cell_font_pt(cell, default=10.0):
    """Taille de police du premier run de la cellule, sinon `default`."""
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                return r.font.size.pt
    return default


def _estimate_row_height_emu(row, col_widths, default_font=10.0):
    """Hauteur rendue d'une ligne = max(hauteur déclarée, besoin du texte wrappé)."""
    needed = row.height
    for ci, cell in enumerate(row.cells):
        if ci >= len(col_widths):
            break
        try:
            ml = cell.margin_left or 0
            mr = cell.margin_right or 0
            mt = cell.margin_top or 0
            mb = cell.margin_bottom or 0
        except Exception:
            ml = mr = 91440
            mt = mb = 45720
        usable_in = max(0.3, (col_widths[ci] - ml - mr) / EMU_PER_IN)
        font_pt = _cell_font_pt(cell, default_font)
        char_w_in = font_pt * CHAR_W_RATIO / 72.0
        chars_per_line = max(1, int(usable_in / char_w_in))
        n_lines = 0
        for p in cell.text_frame.paragraphs:
            txt = "".join(r.text for r in p.runs)
            n_lines += max(1, math.ceil(len(txt) / chars_per_line))
        n_lines = max(1, n_lines)
        line_h_in = font_pt * LINE_SPACING / 72.0
        cell_h = int((n_lines * line_h_in) * EMU_PER_IN) + mt + mb
        needed = max(needed, cell_h)
    return needed


def estimate_table_height_emu(shape, default_font=10.0):
    """Hauteur rendue estimée d'une table (somme des hauteurs de lignes rendues)."""
    table = shape.table
    col_widths = [c.width for c in table.columns]
    return sum(_estimate_row_height_emu(row, col_widths, default_font) for row in table.rows)


def _set_data_rows_font(table, size_pt, header_rows=1):
    """Force la taille de police de toutes les lignes de données."""
    for row in list(table.rows)[header_rows:]:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size_pt)


def fit_table_to_slide(shape, slide_height, *, margin_in=0.15, header_rows=1,
                       font_steps=(9.0, 8.0), default_font=10.0, min_top_in=0.85):
    """
    Garantit qu'une table ne déborde pas du bas de la slide.
    1. Estime la hauteur rendue ; si OK → ne touche à rien.
    2. Sinon réduit la police des lignes de données (font_steps), étape par étape.
    3. Si ça ne suffit pas, remonte la table (sans passer au-dessus de min_top_in).
    Retourne un dict décrivant l'action effectuée (pour log).
    """
    limit = slide_height - int(margin_in * EMU_PER_IN)
    rendered = estimate_table_height_emu(shape, default_font)
    if shape.top + rendered <= limit:
        return {"shape": shape.name, "action": "ok"}

    action = {"shape": shape.name, "action": "font", "font": None, "moved_in": 0.0}
    for step in font_steps:
        _set_data_rows_font(shape.table, step, header_rows)
        action["font"] = step
        rendered = estimate_table_height_emu(shape, step)
        if shape.top + rendered <= limit:
            return action

    # Toujours trop haut : on remonte la table.
    new_top = max(int(min_top_in * EMU_PER_IN), limit - rendered)
    if new_top < shape.top:
        action["moved_in"] = round((shape.top - new_top) / EMU_PER_IN, 2)
        shape.top = new_top
        action["action"] = "font+move"
    return action


def fit_all_tables(prs, *, margin_in=0.15, header_rows=1, verbose=True):
    """
    Applique fit_table_to_slide à toutes les tables de la présentation.
    À appeler systématiquement APRÈS toute injection de texte et AVANT prs.save().
    """
    reports = []
    for idx, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_table:
                continue
            rep = fit_table_to_slide(sh, prs.slide_height,
                                     margin_in=margin_in, header_rows=header_rows)
            rep["slide"] = idx
            reports.append(rep)
            if verbose and rep["action"] != "ok":
                msg = f"  ! slide {idx} {rep['shape']}: police -> {rep['font']} pt"
                if rep.get("moved_in"):
                    msg += f", remontee de {rep['moved_in']}\""
                print(msg)
    return reports


# ---------- Déduplication zip XML (essentiel après édition) ----------

def dedupe_zip(pptx_path):
    """
    Supprime les entrées en doublon dans le zip PPTX, en gardant la dernière
    occurrence. python-pptx peut générer des doublons après édition d'un template
    où des slides ont été retirées, ce qui empêche l'ouverture par LibreOffice
    et PowerPoint.

    À appeler systématiquement après `prs.save(...)`.
    """
    pptx_path = Path(pptx_path)
    tmp_path = pptx_path.with_suffix(pptx_path.suffix + ".tmp")
    seen = {}
    with zipfile.ZipFile(pptx_path, "r") as zin:
        for name in zin.namelist():
            seen[name] = (zin.getinfo(name), zin.read(name))
    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, (info, data) in seen.items():
            zout.writestr(info, data)
    shutil.move(tmp_path, pptx_path)
    return pptx_path
