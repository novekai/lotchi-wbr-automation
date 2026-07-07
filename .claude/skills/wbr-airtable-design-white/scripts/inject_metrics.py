"""Injecte un metrics.json (KM/AP/Comments/CheckList/ActionPlan) dans un PPTX.

Usage :
    python3 scripts/inject_metrics.py \
        --pptx _runs/BAYONNE_WBR_v5.pptx \
        --metrics metrics.json \
        --out _runs/BAYONNE_WBR_v6_final.pptx

Format attendu du JSON :
{
  "slide_3_checklist": [
    {"action": "...", "lotchi": true, "fever": false, "ongoing": false},
    ...
  ],
  "slide_5_gh": {"km": "...", "ap": "..."},
  "slide_7_top_funnel_meta": {"km": "...", "ap": "..."},
  "slide_8_mid_funnel_lp": {"km": "...", "ap": "..."},
  "slide_9_low_funnel_pp": {"km": "...", "ap": "..."},
  "slide_11_soldout": {"km": "...", "ap": "..."},
  "slide_13_bigpicture": {"km": "...", "ap": "..."},
  "slide_14_reach_comment": "Videos : ... | Statics : ...",
  "slide_15_ugc_comment": "...",
  "slide_16_ade_comment": "...",
  "slide_17_rmkt_comment": "...",
  "slide_18_tourist_fr_comment": "...",
  "slide_18_tourist_en_comment": "...",
  "slide_20_actionplan": [
    {"action": "...", "lotchi": "x", "fever": "", "deadline": "W24"},
    ...
  ]
}
"""
from __future__ import annotations

import argparse
import json
import shutil
import sys
import zipfile
from pathlib import Path
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent))
from pptx_utils import fit_all_tables  # noqa: E402


def find_shape(slide, name):
    return next((sh for sh in slide.shapes if sh.name == name), None)


def set_cell(table, row, col, value, size_pt=None):
    """Écrit `value` en forçant un texte clair (fond noir du deck v2).

    Même correctif que build_wbr.set_cell : après un aller-retour Google Slides,
    les runs vides sont supprimés ; on réécrit donc un run blanc, taille lue depuis
    le run existant ou l'endParaRPr de la cellule (défaut 10 pt), pour éviter le
    texte noir invisible sur cellule noire.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    cell = table.cell(row, col)
    tf = cell.text_frame
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    size = None
    for r in p.runs:
        if r.font.size is not None:
            size = r.font.size
            break
    if size is None:
        epr = p._p.find(qn("a:endParaRPr"))
        if epr is not None and epr.get("sz"):
            size = Pt(int(epr.get("sz")) / 100)
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = str(value)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.size = Pt(size_pt) if size_pt else (size or Pt(10))


def add_table_row(table, row_idx, values):
    """Remplit une ligne existante du tableau."""
    for col_idx, v in enumerate(values):
        if col_idx < len(table.columns):
            set_cell(table, row_idx, col_idx, v)


def inject_km_ap(slide, shape_name, km, ap):
    """Injecte Key Metric + Action Plan dans la table 2 cols (header + 1 data row)."""
    sh = find_shape(slide, shape_name)
    if sh is None:
        print(f"  ! {shape_name} introuvable")
        return
    table = sh.table
    # Harmonisation : si la colonne Key Metric a une cellule vide sous la data row
    # (alors que Action Plan est deja fusionnee en rowSpan), fusionner verticalement
    # les deux cellules de la colonne Key Metric pour une hauteur symetrique.
    try:
        if len(table.rows) >= 3:
            top = table.cell(1, 0)
            if not top.is_merge_origin and not top.is_spanned:
                top.merge(table.cell(2, 0))
    except Exception as e:
        print(f"  ! fusion cellule Key Metric ({shape_name}) : {e}")
    set_cell(table, 1, 0, km or "")
    set_cell(table, 1, 1, ap or "")


def _fit_comment_font(slide, sh, tbl, col_idx, text, *, pt_max=10.0, pt_min=8.0):
    """Choisit la plus grande police (entre pt_min et pt_max) qui fait tenir tout le
    texte dans la cellule Comments (largeur colonne x hauteur dispo avant la forme du
    dessous), sans rien tronquer. Troncature propre seulement si, meme a pt_min, ca
    deborde. Renvoie (texte, taille_pt).
    Calibre sur le rendu reel (substituts DejaVu/Carlito) : ~0,56 pt d'avance moyenne
    par caractere, interligne ~1,40.
    """
    import math
    EMU = 914400.0
    txt = text or ""
    try:
        col_w_in = tbl.columns[col_idx].width / EMU
        data_top = sh.top + tbl.rows[0].height
        tops = [s.top for s in slide.shapes
                if s is not sh and s.top is not None and s.top > data_top + int(0.05 * EMU)]
        next_top = min(tops) if tops else (sh.top + sh.height)
        avail_in = (next_top - data_top) / EMU
    except Exception as e:
        print(f"  ! calc police comments : {e}")
        return txt, pt_max
    pt = pt_max
    while pt >= pt_min - 1e-9:
        cpl = max(12.0, (col_w_in * 72.0) / (pt * 0.56))
        line_in = pt * 1.40 / 72.0
        avail_lines = max(1, int(avail_in / line_in))
        need_lines = math.ceil(len(txt) / cpl) if txt else 0
        if need_lines <= avail_lines:
            return txt, round(pt, 1)
        pt -= 0.5
    # deborde meme a pt_min : on tronque au budget de pt_min
    cpl = max(12.0, (col_w_in * 72.0) / (pt_min * 0.56))
    line_in = pt_min * 1.40 / 72.0
    avail_lines = max(1, int(avail_in / line_in))
    budget = int(avail_lines * cpl)
    cut = txt[:budget].rsplit(" ", 1)[0].rstrip(" ,;:.") + "\u2026"
    return cut, pt_min


def inject_comments(slide, shape_name, col_idx, text):
    """Injecte le texte dans la derniere colonne (Comments) de la 1ere data row.
    Reduit automatiquement la police pour tout afficher sans deborder derriere le graphe."""
    sh = find_shape(slide, shape_name)
    if sh is None:
        print(f"  ! {shape_name} introuvable")
        return
    fitted, pt = _fit_comment_font(slide, sh, sh.table, col_idx, text)
    if text:
        if fitted != text:
            print(f"    (comments {shape_name} tronque a pt_min : {len(text)} -> {len(fitted)} c)")
        elif pt < 10.0:
            print(f"    (comments {shape_name} police reduite a {pt} pt)")
    set_cell(sh.table, 1, col_idx, fitted, size_pt=pt)

def inject_checklist(slide, items):
    """Remplit slide 3 checklist : col 0=action, col 1=Done by Lotchi (x si true),
    col 2=Done by Fever, col 3=Ongoing.
    """
    sh = find_shape(slide, "checklist_w1_table")
    if sh is None:
        print("  ! checklist_w1_table introuvable")
        return
    table = sh.table
    # rows[0] = header. Data rows = rows[1:]
    n_rows = len(table.rows)
    for i, item in enumerate(items):
        if i + 1 >= n_rows:
            break  # template limité à N data rows
        set_cell(table, i + 1, 0, item.get("action", ""))
        set_cell(table, i + 1, 1, "x" if item.get("lotchi") else "")
        set_cell(table, i + 1, 2, "x" if item.get("fever") else "")
        set_cell(table, i + 1, 3, "x" if item.get("ongoing") else "")


def inject_actionplan_slide20(slide, items):
    """Remplit la table actionplan slide 20 : 4 cols = Actions / Lotchi / Fever / Deadline."""
    sh = find_shape(slide, "actionplan_table")
    if sh is None:
        print("  ! actionplan_table introuvable")
        return
    table = sh.table
    n_rows = len(table.rows)
    for i, item in enumerate(items):
        if i + 1 >= n_rows:
            break
        set_cell(table, i + 1, 0, item.get("action", ""))
        set_cell(table, i + 1, 1, item.get("lotchi", ""))
        set_cell(table, i + 1, 2, item.get("fever", ""))
        set_cell(table, i + 1, 3, item.get("deadline", ""))


def dedupe_zip_xml(pptx_path: Path):
    tmp = pptx_path.with_suffix(".dedupe.tmp")
    with zipfile.ZipFile(pptx_path, "r") as zin:
        seen = set()
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                if info.filename in seen:
                    continue
                seen.add(info.filename)
                zout.writestr(info, zin.read(info.filename))
    shutil.move(str(tmp), str(pptx_path))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True, type=Path)
    ap.add_argument("--metrics", required=True, type=Path)
    ap.add_argument("--out", required=True, type=Path)
    args = ap.parse_args()

    metrics = json.loads(args.metrics.read_text(encoding="utf-8"))
    prs = Presentation(str(args.pptx))

    # Indexation par slide
    slide3 = prs.slides[2]
    slide5 = prs.slides[4]
    slide7 = prs.slides[6]
    slide8 = prs.slides[7]
    slide9 = prs.slides[8]
    slide11 = prs.slides[10]
    slide13 = prs.slides[12]
    slide14 = prs.slides[13]
    slide15 = prs.slides[14]
    slide16 = prs.slides[15]
    slide17 = prs.slides[16]
    slide18 = prs.slides[17]
    slide20 = prs.slides[19]

    # Check List W-1 (slide 3)
    if "slide_3_checklist" in metrics:
        inject_checklist(slide3, metrics["slide_3_checklist"])
        print("✓ slide 3 checklist injectée")

    # Key Metric / Action Plan tables
    km_ap_mapping = [
        ("slide_5_gh", slide5, "km_ap_table_s5"),
        ("slide_7_top_funnel_meta", slide7, "km_ap_table_s7"),
        ("slide_8_mid_funnel_lp", slide8, "km_ap_table_s8"),
        ("slide_9_low_funnel_pp", slide9, "km_ap_table_s9"),
        ("slide_11_soldout", slide11, "km_ap_table_s11"),
        ("slide_13_bigpicture", slide13, "km_ap_table_s13"),
    ]
    for key, slide, shape_name in km_ap_mapping:
        if key in metrics:
            inject_km_ap(slide, shape_name, metrics[key].get("km"), metrics[key].get("ap"))
            print(f"✓ {key} KM/AP injecté")

    # Comments par cluster (slides 14-18)
    # Pour REACH (slide 14) : col 2 = comments (la table a 3 cols : Campaign, CTRs, Comments)
    if "slide_14_reach_comment" in metrics:
        inject_comments(slide14, "reach_comments_table", 2, metrics["slide_14_reach_comment"])
        print("✓ slide 14 REACH comment")
    # Pour UGC/ADE/RMKT (slides 15/16/17) : col 3 = comments (table 4 cols : Campaign, CTRs, ROAS, Comments)
    for key, slide, shape_name in [
        ("slide_15_ugc_comment", slide15, "ugc_comments_table"),
        ("slide_16_ade_comment", slide16, "ade_comments_table"),
        ("slide_17_rmkt_comment", slide17, "rmkt_comments_table"),
        ("slide_18_tourist_fr_comment", slide18, "tourist_fr_comments_table"),
        ("slide_18_tourist_en_comment", slide18, "tourist_en_comments_table"),
    ]:
        if key in metrics:
            inject_comments(slide, shape_name, 3, metrics[key])
            print(f"✓ {key}")

    # Action Plan slide 20
    if "slide_20_actionplan" in metrics:
        inject_actionplan_slide20(slide20, metrics["slide_20_actionplan"])
        print(f"✓ slide 20 action plan ({len(metrics['slide_20_actionplan'])} actions)")

    # Garde-fou anti-débordement : aucune table ne doit sortir du bas de slide
    print("Contrôle débordement des tables…")
    fit_all_tables(prs)

    # Sauvegarde
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    dedupe_zip_xml(args.out)
    print(f"\nOK : {args.out}")


if __name__ == "__main__":
    main()
