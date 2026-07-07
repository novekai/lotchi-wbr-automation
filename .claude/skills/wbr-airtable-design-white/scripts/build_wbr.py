"""
Orchestrateur du skill `wbr-airtable` — produit le WBR Full v2 d'une ville
à partir de la base Airtable Lotchi WBR.

Usage :
    python3 scripts/build_wbr.py --city Bayonne --week 2026-W23 \
        --template wbr_template.pptx --out output/BAYONNE_WBR.pptx

Le script :
1. Charge les données via `airtable_fetch.WBRDataLoader`.
2. Classe les créatives via `cluster_classifier`.
3. Calcule les agrégats (KPI summary, Big Picture CTR/CPM).
4. Génère les charts matplotlib via `generate_charts`.
5. Édite le template Full v2 via `pptx_utils` :
   - placeholders {{VILLE_UPPER}} / {{ANNEE}} / {{DATE_MERCREDI}}
   - tables KPI summary (slides 14, 15) et Big Picture (slide 13)
   - 3 tables sold-out (slide 11)
   - 6 tableaux de créas par cluster (slides 14-18)
   - charts slides 5, 7, 8, 9
6. Laisse vides les Key Metric / Action Plan / Comments tables, qui sont
   rédigés par Claude post-build via le playbook de décisions.
7. Sauvegarde le PPTX final avec déduplication zip XML.

Le script peut être invoqué de bout en bout, ou en deux étapes via
`--stage extract` puis `--stage render` pour permettre une rédaction
manuelle des Key Metrics entre les deux.
"""
from __future__ import annotations

import argparse
import datetime as dt
import json
import os
import shutil
import sys
import tempfile
import time
import zipfile
from collections import defaultdict
from pathlib import Path

# Ajouter le dossier scripts au path pour les imports relatifs
sys.path.insert(0, str(Path(__file__).resolve().parent))

from airtable_fetch import AirtableClient, WBRDataLoader  # noqa: E402
from cluster_classifier import (  # noqa: E402
    DISPLAY_CLUSTERS,
    aggregate_cluster_kpis,
    classify_creative,
    group_by_cluster,
)
import generate_charts as gc  # noqa: E402
from time_variables import compute_dates  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402


# Ratios W/H des zones assets du template v2 verrouille (evite toute distorsion
# a l'injection : l'image est rendue au ratio exact de la boite cible).
ASSET_ASPECT = {
    "REACH": 4.366, "UGC": 4.366, "ADE": 3.092, "RMKT": 3.092,
    "TOURIST_FR": 7.603, "TOURIST_EN": 7.603, "TOURIST_ES": 7.603,
}


# ---------------------------------------------------------------------------
# Helpers de traitement de données


def aggregate_marketing_channels(records: list[dict]) -> dict:
    """Indexe les enregistrements Marketing Channel Performance.

    Retourne : {week_id: {channel_name: {ctr, cpm, impressions, spend, roas}}}.
    """
    out: dict[str, dict[str, dict]] = defaultdict(dict)
    # Pour résoudre les noms des canaux et semaines, on garde un index par
    # record_id Airtable, qu'on reconstitue via les autres dumps.
    for r in records:
        f = r.get("fields", {})
        key = f.get("Channel perf key", "")
        # Channel perf key = "Show ID | Week ID | Channel name"
        parts = [p.strip() for p in key.split("|")]
        if len(parts) != 3:
            continue
        _show_id, week_id, channel_name = parts
        out[week_id][channel_name] = {
            "impressions": int(f.get("Impressions") or 0),
            "clicks": int(f.get("Clicks") or 0),
            "ctr": float(f.get("CTR") or 0.0),
            "spend": float(f.get("Spend") or 0.0),
            "cpm": float(f.get("CPM") or 0.0),
            "roas_cumul": float(f.get("ROAS cumul") or 0.0),
            "roas_weekly": float(f.get("ROAS weekly") or 0.0),
            "conversions": int(f.get("Conversions") or 0),
        }
    return dict(out)


def aggregate_landing_page(records: list[dict]) -> dict:
    """Indexe : {week_id: {channel_group: {visitors, ctr, cta_clicks, page_views}}}."""
    out: dict[str, dict[str, dict]] = defaultdict(dict)
    for r in records:
        f = r.get("fields", {})
        key = f.get("LP key", "")
        parts = [p.strip() for p in key.split("|")]
        if len(parts) != 3:
            continue
        _show_id, week_id, channel_group = parts
        out[week_id][channel_group] = {
            "visitors": int(f.get("Visitors") or 0),
            "page_views": int(f.get("Page views") or 0),
            "cta_clicks": int(f.get("CTA clicks") or 0),
            "ctr": float(f.get("CTR") or 0.0),
        }
    return dict(out)


def aggregate_weekly_kpis(records: list[dict]) -> dict:
    """Indexe : {week_id: champs}."""
    out: dict[str, dict] = {}
    for r in records:
        f = r.get("fields", {})
        key = f.get("KPI key", "")
        parts = [p.strip() for p in key.split("|")]
        if len(parts) != 2:
            continue
        _show_id, week_id = parts
        out[week_id] = f
    return out


def aggregate_funnel(records: list[dict]) -> dict:
    """Indexe : {week_id: champs}."""
    out: dict[str, dict] = {}
    for r in records:
        f = r.get("fields", {})
        key = f.get("Funnel key", "")
        parts = [p.strip() for p in key.split("|")]
        if len(parts) != 2:
            continue
        _show_id, week_id = parts
        out[week_id] = f
    return out


def aggregate_sessions(records: list[dict]) -> dict:
    """Indexe : {week_id: [session_dict, ...]} prêt pour soldout_table."""
    out: dict[str, list[dict]] = defaultdict(list)
    for r in records:
        f = r.get("fields", {})
        key = f.get("Session key", "")
        # Session key = "Show ID | Event date | Event time"
        parts = [p.strip() for p in key.split("|")]
        if len(parts) != 3:
            continue
        event_date = f.get("Event date") or parts[1]
        # Déduire la semaine ISO
        try:
            y, w, _ = dt.date.fromisoformat(event_date).isocalendar()
            week_id = f"{y}-W{w:02d}"
        except Exception:
            continue
        out[week_id].append({
            "date": event_date,
            "time": f.get("Event time") or parts[2],
            "cap": int(f.get("Capacity") or 0),
            "sold": int(f.get("Tickets sold") or 0),
            "fill": float(f.get("Fill rate") or 0.0),
        })
    return dict(out)


def extract_creatives_payload(records: list[dict]) -> list[dict]:
    """Réduit les enregistrements créa à un payload léger pour les charts/tables."""
    out = []
    for r in records:
        f = r.get("fields", {})
        out.append({
            "creative_id": f.get("Creative ID"),
            "creative_name": f.get("Creative name"),
            "format": f.get("Format"),
            "preview_url": f.get("Preview URL"),
            "impressions": int(f.get("Impressions") or 0),
            "clicks": int(f.get("Clicks") or 0),
            "ctr": float(f.get("CTR") or 0.0),
            "spend": float(f.get("Spend") or 0.0),
            "roas": float(f.get("ROAS") or 0.0),
            "roas_utm": float(f.get("ROAS UTM") or 0.0),  # à venir côté Airtable
            "cpm": float(f.get("CPM") or 0.0),
            "cpc": float(f.get("CPC") or 0.0),
            "date_posted": f.get("Date posted"),
            "objective": f.get("Campaign Objective"),
            "adset": f.get("Ad set name"),
            "phase": f.get("Phase WL/SL"),
            "cluster": classify_creative(f),
        })
    return out


def aggregate_creatives_subset(creatives_list: list[dict], cluster: str) -> dict:
    """Agrégat pondéré par impressions sur les créas d'un cluster donné.

    creatives_list : payload léger (post extract_creatives_payload).
    cluster : 'REACH', 'UGC', 'ADE', 'RMKT', 'TOURIST_FR', 'TOURIST_EN', 'TOURIST_ES'...

    Le ROAS effectif est `roas_utm` si dispo (champ Datacloud), sinon `roas`
    (champ Fever Zone, vide pour les snapshots weekly_datacloud).
    """
    imps_total = clicks_total = 0
    spend_total = revenue_total = 0.0
    for c in creatives_list:
        if c["cluster"] != cluster:
            continue
        imps_total += c["impressions"]
        clicks_total += c["clicks"]
        spend_total += c["spend"]
        roas_eff = c.get("roas_utm") or c.get("roas") or 0.0
        if roas_eff:
            revenue_total += c["spend"] * roas_eff
    ctr = (clicks_total / imps_total) if imps_total else 0.0
    cpm = (spend_total / imps_total * 1000) if imps_total else 0.0
    roas = (revenue_total / spend_total) if spend_total else 0.0
    return {
        "impressions": imps_total, "clicks": clicks_total, "spend": spend_total,
        "ctr": ctr, "cpm": cpm, "roas": roas,
    }


# ---------------------------------------------------------------------------
# Édition du PPTX


def find_shape(slide, name):
    return next((sh for sh in slide.shapes if sh.name == name), None)


def replace_placeholders(prs, mapping: dict[str, str]):
    """Remplace toute occurrence de {{KEY}} dans les shapes texte."""
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    txt = r.text
                    for k, v in mapping.items():
                        token = f"{{{{{k}}}}}"
                        if token in txt:
                            txt = txt.replace(token, v)
                    if txt != r.text:
                        r.text = txt


def set_cell(table, row, col, value):
    """Écrit `value` dans une cellule de données en forçant un texte clair.

    Les cellules de données du template v2 (fond noir) doivent afficher un texte
    clair. Après un aller-retour Google Slides, les runs vides sont supprimés :
    on ne peut pas se fier au formatage existant. On (ré)écrit donc un run avec
    couleur claire (#E6E6E6) et une taille par défaut si aucune n'est présente,
    pour éviter le cas « texte noir sur cellule noire » (invisible).
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    cell = table.cell(row, col)
    tf = cell.text_frame
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    size = None
    for r in p.runs:
        if r.font.size is not None:
            size = r.font.size
            break
    if size is None:
        epr = p._p.find(qn("a:endParaRPr"))
        if epr is not None and epr.get("sz"):
            size = Pt(int(epr.get("sz")) / 100)
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = str(value)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.size = size or Pt(10)


def replace_picture(slide, shape_name, image_path):
    """Remplace l'image d'un shape image existant en conservant ses dimensions."""
    sh = find_shape(slide, shape_name)
    if sh is None:
        print(f"  ! shape {shape_name!r} absent", file=sys.stderr)
        return
    L, T, W, H = sh.left, sh.top, sh.width, sh.height
    sh._element.getparent().remove(sh._element)
    slide.shapes.add_picture(str(image_path), L, T, width=W, height=H)


def dedupe_zip_xml(pptx_path: Path):
    """Workaround python-pptx : retire les entrées dupliquées dans le zip."""
    tmp = pptx_path.with_suffix(".dedupe.tmp")
    with zipfile.ZipFile(pptx_path, "r") as zin:
        seen = set()
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                if info.filename in seen:
                    continue
                seen.add(info.filename)
                zout.writestr(info, zin.read(info.filename))
    shutil.move(str(tmp), str(pptx_path))


# ---------------------------------------------------------------------------
# Pipeline principal


def build(city: str, week_id: str, template_path: Path, out_path: Path,
          tmp_dir: Path | None = None, data: dict | None = None):
    t0 = time.time()
    tmp_dir = Path(tmp_dir or tempfile.mkdtemp(prefix="wbr_"))
    tmp_dir.mkdir(parents=True, exist_ok=True)
    print(f"[build] tmp_dir = {tmp_dir}")

    # 1. Dates
    dates = compute_dates(dt.date.today())
    # Si l'utilisateur a passé une semaine explicite, on s'aligne dessus
    # (en supposant que week_id correspond à dates.week_id_minus1, c'est-à-dire
    # la semaine analysée).
    # Pour le moment, on accepte la divergence et on utilise week_id pour
    # filtrer Airtable, mais on calque date_mercredi sur la convention.
    iso_year_w, iso_week_w = (int(week_id[:4]), int(week_id[-2:]))
    lundi_w = dt.date.fromisocalendar(iso_year_w, iso_week_w, 1)
    dimanche_w = dt.date.fromisocalendar(iso_year_w, iso_week_w, 7)
    # Mercredi DE la semaine passée en --week (jour ISO 3), pas de la suivante.
    mercredi_pres = dt.date.fromisocalendar(iso_year_w, iso_week_w, 3)
    date_mercredi = mercredi_pres.strftime("%d/%m/%Y")

    print(f"[build] city={city} week={week_id} mercredi_pres={date_mercredi}")

    # 2. Chargement Airtable (ou data injectee pour test hors-ligne)
    if data is None:
        print("[build] Airtable fetch...")
        client = AirtableClient()
        loader = WBRDataLoader(client, city=city, week_id=week_id)
        data = loader.load_all()
    else:
        print("[build] Data injectee (test hors Airtable)")
    print(f"[build]   show={data['show']['fields'].get('Show name')}")
    print(f"[build]   week_ids={data['week_ids']}")
    print(f"[build]   sessions={len(data['sessions'])}, kpis={len(data['weekly_kpis'])}, "
          f"channels={len(data['marketing_channels'])}, "
          f"lp={len(data['landing_page'])}, funnel={len(data['funnel'])}, "
          f"creatives_w1={len(data['creatives_wminus1'])}, "
          f"creatives_w2={len(data['creatives_wminus2'])}")

    # 3. Agrégations
    weekly_kpis = aggregate_weekly_kpis(data["weekly_kpis"])
    marketing = aggregate_marketing_channels(data["marketing_channels"])
    landing = aggregate_landing_page(data["landing_page"])
    funnel = aggregate_funnel(data["funnel"])
    sessions_by_week = aggregate_sessions(data["sessions"])
    creatives_wminus1 = extract_creatives_payload(data["creatives_wminus1"])
    creatives_wminus2 = extract_creatives_payload(data["creatives_wminus2"])

    # Semaines de référence dans la sémantique du WBR :
    # [W-3, W-2, W-1, Current, W+1] où W-1 = semaine analysée
    w_minus3, w_minus2, w_minus1, w_curr, w_plus1 = data["week_ids"]

    # 4. Génération des charts
    print("[build] Génération charts...")
    out_charts = tmp_dir / "charts"
    out_charts.mkdir(exist_ok=True)

    # Charts slides 5-9 — 3 dernières semaines analysées : W-3, W-2, W-1
    weeks_3 = [w_minus3, w_minus2, w_minus1]
    revenue_3 = [float(weekly_kpis.get(w, {}).get("Total revenue") or 0) for w in weeks_3]
    spend_3 = [float(weekly_kpis.get(w, {}).get("Mkt spend") or 0) for w in weeks_3]
    roas_3 = [float(weekly_kpis.get(w, {}).get("Avg ROAS") or 0) for w in weeks_3]
    # Acc ROI (ROI cumule sur la periode analysee) : revenu cumule / spend cumule
    acc_roas_3 = []
    _cr = _cs = 0.0
    for _rv, _sp in zip(revenue_3, spend_3):
        _cr += _rv; _cs += _sp
        acc_roas_3.append((_cr / _cs) if _cs else 0.0)
    chart_mkt = out_charts / "chart_mkt_performance.png"
    gc.chart_mkt_performance(weeks_3, revenue_3, spend_3, roas_3, out_path=chart_mkt,
                             acc_roas=acc_roas_3)

    # Chart Slide 5 bas (tickets sold)
    tickets_3 = [int(weekly_kpis.get(w, {}).get("Total tickets") or 0) for w in weeks_3]
    chart_tickets = out_charts / "chart_tickets_sold.png"
    gc.chart_tickets_sold(weeks_3, tickets_3, chart_tickets)

    # Chart Slide 7 top funnel meta — CTR Meta et Google 3 semaines.
    # Les canaux réels dans Airtable sont : meta_fb_ig, google_search,
    # google_display, google_youtube, google_demand_gen.
    # On agrège tous les google_* (pondéré par impressions) en un seul "Google".
    META_CHANNEL = "meta_fb_ig"
    GOOGLE_CHANNELS = ["google_search", "google_display",
                       "google_youtube", "google_demand_gen"]

    def channel_metric_agg(week: str, channels: list[str], metric: str) -> float:
        """Agrégation pondérée par impressions sur une liste de canaux.

        Pour CTR : moyenne pondérée. Pour CPM : moyenne pondérée par impressions.
        Pour ROAS : moyenne pondérée par spend.
        Pour impressions/clicks/spend : somme.
        """
        ch_data = [marketing.get(week, {}).get(c) for c in channels]
        ch_data = [c for c in ch_data if c]
        if not ch_data:
            return 0.0
        if metric == "impressions" or metric == "clicks" or metric == "spend":
            return sum(c.get(metric, 0) for c in ch_data)
        if metric == "ctr":
            imps_total = sum(c.get("impressions", 0) for c in ch_data)
            if not imps_total:
                return 0.0
            return sum(c.get("ctr", 0) * c.get("impressions", 0) for c in ch_data) / imps_total
        if metric == "cpm":
            imps_total = sum(c.get("impressions", 0) for c in ch_data)
            if not imps_total:
                return 0.0
            return sum(c.get("cpm", 0) * c.get("impressions", 0) for c in ch_data) / imps_total
        if metric == "roas":
            spend_total = sum(c.get("spend", 0) for c in ch_data)
            if not spend_total:
                return 0.0
            return sum(c.get("roas_cumul", 0) * c.get("spend", 0) for c in ch_data) / spend_total
        return 0.0

    def channel_ctr(week, channels):
        return channel_metric_agg(week, channels, "ctr") * 100  # en %

    meta_ctr = [channel_ctr(w, [META_CHANNEL]) for w in weeks_3]
    google_ctr = [channel_ctr(w, GOOGLE_CHANNELS) for w in weeks_3]
    chart_meta = out_charts / "chart_top_funnel_meta.png"
    gc.chart_top_funnel_ctrs(weeks_3, meta_ctr, google_ctr, chart_meta)

    # Chart Slide 8 landing page
    LP_TOTAL_CHANNEL = "(Tous canaux — total Fever)"
    LP_FALLBACK_CHANNEL = "fever"

    def lp_total(week, field):
        # Priorité à la ligne agrégée « Tous canaux — total Fever » ; si elle
        # n'existe pas pour cette ville/semaine (cas Reims), on retombe sur la
        # ligne « fever » qui porte les mêmes Visitors/CTA/CTR côté Fever.
        week_data = landing.get(week, {})
        all_ch = week_data.get(LP_TOTAL_CHANNEL)
        if not all_ch:
            all_ch = week_data.get(LP_FALLBACK_CHANNEL, {})
        return all_ch.get(field, 0)
    lp_visitors = [lp_total(w, "visitors") for w in weeks_3]
    lp_cta = [lp_total(w, "cta_clicks") for w in weeks_3]
    lp_ctr = [lp_total(w, "ctr") * 100 for w in weeks_3]
    chart_lp = out_charts / "chart_landing_page.png"
    gc.chart_landing_page(weeks_3, lp_visitors, lp_cta, lp_ctr, chart_lp)

    # Chart Slide 9 product page
    pp_views = [int(funnel.get(w, {}).get("Product page views") or 0) for w in weeks_3]
    pp_cr = [float(funnel.get(w, {}).get("Estimated conversion rate") or 0) * 100
             for w in weeks_3]
    chart_pp = out_charts / "chart_product_page.png"
    gc.chart_product_page(weeks_3, pp_views, pp_cr, chart_pp)

    # Tables sold-out W-1, W, W+1
    def make_soldout(week, name):
        path = out_charts / f"soldout_{name}.png"
        sessions = sessions_by_week.get(week, [])
        if not sessions:
            gc.soldout_placeholder(
                f"Pas de séance enregistrée pour {week}.", path
            )
        else:
            gc.soldout_table(sessions, path)
        return path
    soldout_wm1 = make_soldout(w_minus1, "wminus1")
    soldout_curr = make_soldout(w_curr, "current")
    soldout_wp1 = make_soldout(w_plus1, "wplus1")

    # Tableaux créa par cluster (slides 14-18). Source : créas snapshot W-1.
    assets_paths = {}
    for cluster in DISPLAY_CLUSTERS:
        cluster_creatives = [c for c in creatives_wminus1 if c["cluster"] == cluster]
        path = out_charts / f"assets_{cluster.lower()}.png"
        gc.creative_assets_table(cluster_creatives, path, aspect=ASSET_ASPECT.get(cluster))
        assets_paths[cluster] = path

    # KPI summary slides 14, 15 — vraies agrégations W-2 vs W-1 via les
    # 2 snapshots distincts.
    def kpis_cluster_w2_w1(cluster):
        return {
            "wminus2": aggregate_creatives_subset(creatives_wminus2, cluster),
            "wminus1": aggregate_creatives_subset(creatives_wminus1, cluster),
        }

    reach_kpis = kpis_cluster_w2_w1("REACH")
    ugc_kpis = kpis_cluster_w2_w1("UGC")
    # Comments tables slides 14-18 : agrégat W-1 du cluster pour remplir
    # les cellules CTR (col 1) et ROAS (col 2 quand présente)
    cluster_w1_agg = {
        cluster: aggregate_creatives_subset(creatives_wminus1, cluster)
        for cluster in DISPLAY_CLUSTERS
    }

    # Big Picture (slide 13) — CTR / CPM par canal sur W-3, W-2, W-1
    # + colonnes Evol (W-2 vs W-3, puis W-1 vs W-2).
    def evol_pct(curr: float, prev: float) -> str:
        if not prev:
            return "—"
        v = (curr - prev) / prev * 100
        sign = "+" if v >= 0 else ""
        return f"{sign}{v:.2f}%"

    bigpic = {}
    for label, channels in (("Meta", [META_CHANNEL]), ("Google", GOOGLE_CHANNELS)):
        ctr_vals = {w: channel_metric_agg(w, channels, "ctr") * 100
                    for w in [w_minus3, w_minus2, w_minus1]}
        cpm_vals = {w: channel_metric_agg(w, channels, "cpm")
                    for w in [w_minus3, w_minus2, w_minus1]}
        bigpic[label] = {
            "ctr": ctr_vals,
            "cpm": cpm_vals,
            "ctr_evol_w2_w3": evol_pct(ctr_vals[w_minus2], ctr_vals[w_minus3]),
            "ctr_evol_w1_w2": evol_pct(ctr_vals[w_minus1], ctr_vals[w_minus2]),
            "cpm_evol_w2_w3": evol_pct(cpm_vals[w_minus2], cpm_vals[w_minus3]),
            "cpm_evol_w1_w2": evol_pct(cpm_vals[w_minus1], cpm_vals[w_minus2]),
        }

    # 5. Édition du template
    print("[build] Édition template PPTX...")
    prs = Presentation(str(template_path))

    # Placeholders globaux
    ville_upper = city.upper()
    annee = str(iso_year_w)
    replace_placeholders(prs, {
        "VILLE_UPPER": ville_upper,
        "VILLE": city,
        "ANNEE": annee,
        "DATE_MERCREDI": date_mercredi,
    })

    # Slide 5 — charts
    slide5 = prs.slides[4]
    replace_picture(slide5, "chart_mkt_performance", chart_mkt)
    replace_picture(slide5, "chart_tickets_sold", chart_tickets)

    # Slide 7 — chart meta
    slide7 = prs.slides[6]
    replace_picture(slide7, "chart_top_funnel_meta", chart_meta)

    # Slide 8 — chart LP
    slide8 = prs.slides[7]
    replace_picture(slide8, "chart_landing_page", chart_lp)

    # Slide 9 — chart PP
    slide9 = prs.slides[8]
    replace_picture(slide9, "chart_product_page", chart_pp)

    # Slide 11 — soldout
    slide11 = prs.slides[10]
    replace_picture(slide11, "soldout_wminus1", soldout_wm1)
    replace_picture(slide11, "soldout_current", soldout_curr)
    replace_picture(slide11, "soldout_wplus1", soldout_wp1)

    # Slide 13 — Big Picture CTR + CPM avec colonnes Evol
    # Header table (7 cols) : [Métrique, W-3, Evol, W-2, Evol, W-1, Bench]
    slide13 = prs.slides[12]
    ctr_table = find_shape(slide13, "bigpicture_ctr_table").table
    cpm_table = find_shape(slide13, "bigpicture_cpm_table").table
    for r_idx, channel in [(1, "Meta"), (2, "Google")]:
        ctrs = bigpic[channel]["ctr"]
        cpms = bigpic[channel]["cpm"]
        # Valeurs W-3, W-2, W-1
        set_cell(ctr_table, r_idx, 1, f"{ctrs[w_minus3]:.2f}%")
        set_cell(ctr_table, r_idx, 3, f"{ctrs[w_minus2]:.2f}%")
        set_cell(ctr_table, r_idx, 5, f"{ctrs[w_minus1]:.2f}%")
        set_cell(cpm_table, r_idx, 1, f"{cpms[w_minus3]:.2f} $")
        set_cell(cpm_table, r_idx, 3, f"{cpms[w_minus2]:.2f} $")
        set_cell(cpm_table, r_idx, 5, f"{cpms[w_minus1]:.2f} $")
        # Évolutions (col 2 = W-2 vs W-3, col 4 = W-1 vs W-2)
        set_cell(ctr_table, r_idx, 2, bigpic[channel]["ctr_evol_w2_w3"])
        set_cell(ctr_table, r_idx, 4, bigpic[channel]["ctr_evol_w1_w2"])
        set_cell(cpm_table, r_idx, 2, bigpic[channel]["cpm_evol_w2_w3"])
        set_cell(cpm_table, r_idx, 4, bigpic[channel]["cpm_evol_w1_w2"])

    # Helpers d'affichage : tiret si cluster vide (impressions 0)
    def fmt_imp(v):
        return f"{v:,}".replace(",", " ") if v else "—"
    def fmt_pct(v):
        return f"{v*100:.2f}%" if v else "—"
    def fmt_num(v):
        return f"{v:.2f}" if v else "—"
    def fmt_cur(v):
        return f"{v:.2f} $" if v else "—"

    # Slide 14 — REACH KPI summary + Comments table
    slide14 = prs.slides[13]
    reach_table = find_shape(slide14, "reach_kpi_summary_table").table
    k_m2 = reach_kpis["wminus2"]; k_m1 = reach_kpis["wminus1"]
    set_cell(reach_table, 1, 0, fmt_imp(k_m2['impressions']))
    set_cell(reach_table, 1, 1, fmt_imp(k_m1['impressions']))
    set_cell(reach_table, 1, 2, fmt_pct(k_m2['ctr']))
    set_cell(reach_table, 1, 3, fmt_pct(k_m1['ctr']))
    set_cell(reach_table, 1, 4, fmt_num(k_m2['roas']))
    set_cell(reach_table, 1, 5, fmt_num(k_m1['roas']))
    set_cell(reach_table, 1, 6, fmt_cur(k_m2['cpm']))
    set_cell(reach_table, 1, 7, fmt_cur(k_m1['cpm']))
    # Comments table REACH (cols : Campaign / CTRs / Comments) — REACH n'a
    # pas de col ROAS sur ce template
    reach_comments = find_shape(slide14, "reach_comments_table").table
    set_cell(reach_comments, 1, 1, f"{cluster_w1_agg['REACH']['ctr']*100:.2f}%")
    replace_picture(slide14, "chart_reach_assets", assets_paths["REACH"])

    # Slide 15 — UGC KPI summary + Comments table
    slide15 = prs.slides[14]
    ugc_table = find_shape(slide15, "ugc_kpi_summary_table").table
    k_m2 = ugc_kpis["wminus2"]; k_m1 = ugc_kpis["wminus1"]
    set_cell(ugc_table, 1, 0, fmt_imp(k_m2['impressions']))
    set_cell(ugc_table, 1, 1, fmt_imp(k_m1['impressions']))
    set_cell(ugc_table, 1, 2, fmt_pct(k_m2['ctr']))
    set_cell(ugc_table, 1, 3, fmt_pct(k_m1['ctr']))
    set_cell(ugc_table, 1, 4, fmt_num(k_m2['roas']))
    set_cell(ugc_table, 1, 5, fmt_num(k_m1['roas']))
    set_cell(ugc_table, 1, 6, fmt_cur(k_m2['cpm']))
    set_cell(ugc_table, 1, 7, fmt_cur(k_m1['cpm']))
    # Comments table UGC (Campaign / CTRs / ROAS / Comments)
    ugc_comments = find_shape(slide15, "ugc_comments_table").table
    set_cell(ugc_comments, 1, 1, f"{cluster_w1_agg['UGC']['ctr']*100:.2f}%")
    set_cell(ugc_comments, 1, 2, f"{cluster_w1_agg['UGC']['roas']:.2f}")
    replace_picture(slide15, "chart_ugc_assets", assets_paths["UGC"])

    # Slide 16 — ADE (Comments table seulement)
    slide16 = prs.slides[15]
    ade_comments = find_shape(slide16, "ade_comments_table").table
    set_cell(ade_comments, 1, 1, f"{cluster_w1_agg['ADE']['ctr']*100:.2f}%")
    set_cell(ade_comments, 1, 2, f"{cluster_w1_agg['ADE']['roas']:.2f}")
    replace_picture(slide16, "chart_ade_assets", assets_paths["ADE"])

    # Slide 17 — RMKT (Comments table seulement)
    slide17 = prs.slides[16]
    rmkt_comments = find_shape(slide17, "rmkt_comments_table").table
    set_cell(rmkt_comments, 1, 1, f"{cluster_w1_agg['RMKT']['ctr']*100:.2f}%")
    set_cell(rmkt_comments, 1, 2, f"{cluster_w1_agg['RMKT']['roas']:.2f}")
    replace_picture(slide17, "chart_rmkt_assets", assets_paths["RMKT"])

    # Slide 18 — TOURIST FR + TOURIST EN (séparation par langue)
    slide18 = prs.slides[17]
    tourist_fr_comments = find_shape(slide18, "tourist_fr_comments_table").table
    set_cell(tourist_fr_comments, 1, 1, f"{cluster_w1_agg['TOURIST_FR']['ctr']*100:.2f}%")
    set_cell(tourist_fr_comments, 1, 2, f"{cluster_w1_agg['TOURIST_FR']['roas']:.2f}")
    tourist_en_comments = find_shape(slide18, "tourist_en_comments_table").table
    set_cell(tourist_en_comments, 1, 1, f"{cluster_w1_agg['TOURIST_EN']['ctr']*100:.2f}%")
    set_cell(tourist_en_comments, 1, 2, f"{cluster_w1_agg['TOURIST_EN']['roas']:.2f}")
    replace_picture(slide18, "chart_tourist_fr_assets", assets_paths["TOURIST_FR"])
    replace_picture(slide18, "chart_tourist_en_assets", assets_paths["TOURIST_EN"])

    # 6. Garde-fou anti-débordement + sauvegarde + déduplication
    from pptx_utils import fit_all_tables
    print("[build] Contrôle débordement des tables…")
    fit_all_tables(prs)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    dedupe_zip_xml(out_path)

    duration = time.time() - t0
    print(f"[build] OK en {duration:.1f}s → {out_path}")
    return {
        "out_path": str(out_path),
        "duration_s": duration,
        "city": city,
        "week_id": week_id,
        "show_name": data["show"]["fields"].get("Show name"),
        "counts": {
            "sessions": len(data["sessions"]),
            "weekly_kpis": len(data["weekly_kpis"]),
            "marketing_channels": len(data["marketing_channels"]),
            "landing_page": len(data["landing_page"]),
            "funnel": len(data["funnel"]),
            "creatives_wminus1": len(data["creatives_wminus1"]),
            "creatives_wminus2": len(data["creatives_wminus2"]),
            "creatives_by_cluster_w1": {
                c: sum(1 for x in creatives_wminus1 if x["cluster"] == c)
                for c in DISPLAY_CLUSTERS + ["TOURIST_OTHER", "WL", "OG", "OTHER"]
            },
        },
    }


# ---------------------------------------------------------------------------
# CLI


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--city", required=True)
    parser.add_argument("--week", required=True, help="Semaine ISO, ex 2026-W23")
    parser.add_argument("--template", required=True, type=Path)
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--tmp-dir", type=Path, default=None)
    parser.add_argument("--report-json", type=Path, default=None,
                        help="Si fourni, écrit le rapport de run en JSON.")
    args = parser.parse_args()

    report = build(
        city=args.city,
        week_id=args.week,
        template_path=args.template,
        out_path=args.out,
        tmp_dir=args.tmp_dir,
    )

    if args.report_json:
        args.report_json.parent.mkdir(parents=True, exist_ok=True)
        args.report_json.write_text(
            json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8"
        )
        print(f"[main] report JSON -> {args.report_json}")

    print(f"[main] Termine : {report['out_path']} ({report['duration_s']:.1f}s)")


if __name__ == "__main__":
    main()
